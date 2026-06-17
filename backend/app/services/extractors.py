# app/services/extractors.py
import io
import asyncio
import re
from urllib.parse import urlparse
from typing import List, Dict, Any

import httpx
import trafilatura
from bs4 import BeautifulSoup
from pypdf import PdfReader
from pptx import Presentation
from youtube_transcript_api import YouTubeTranscriptApi


# ---------------------------------------------------------------------------
# PDF Extraction  ||  Library: pypdf
# ---------------------------------------------------------------------------
async def extract_text_from_pdf(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    """
    Reads a PDF from raw bytes, extracts text page by page,
    and attaches citation metadata.
    """
    pdf_file = io.BytesIO(file_bytes)
    reader = PdfReader(pdf_file)
    extracted_pages = []

    for page_num, page in enumerate(reader.pages):
        raw_text = page.extract_text()
        if raw_text:
            clean_text = raw_text.replace('\n', ' ').strip()
            if len(clean_text) > 10:
                extracted_pages.append({
                    "content": clean_text,
                    "metadata": {
                        "source": filename,
                        "page": page_num + 1,
                        "type": "pdf"
                    }
                })

    return extracted_pages


# ---------------------------------------------------------------------------
# PPTX Extraction  ||  Library: python-pptx
# ---------------------------------------------------------------------------
async def extract_text_from_pptx(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    """
    Reads a PPTX from raw bytes, extracts text slide by slide,
    and attaches citation metadata.
    """
    pptx_file = io.BytesIO(file_bytes)
    presentation = Presentation(pptx_file)
    extracted_slides = []

    for slide_num, slide in enumerate(presentation.slides):
        slide_text_elements = []
        for shape in slide.shapes:
            if hasattr(shape, "text") and shape.text:
                slide_text_elements.append(shape.text)

        raw_text = " ".join(slide_text_elements)
        if raw_text:
            clean_text = raw_text.replace('\n', ' ').replace('\x0b', ' ').strip()
            if len(clean_text) > 10:
                extracted_slides.append({
                    "content": clean_text,
                    "metadata": {
                        "source": filename,
                        "slide": slide_num + 1,
                        "type": "pptx"
                    }
                })

    return extracted_slides


# ---------------------------------------------------------------------------
# YouTube Transcript Extraction  ||  Library: youtube-transcript-api
# ---------------------------------------------------------------------------
async def extract_text_from_youtube(url: str) -> List[Dict[str, Any]]:
    """
    Extracts the video ID from a YouTube URL, fetches the transcript,
    and attaches timestamp metadata.
    Runs the blocking API call in a thread executor to avoid blocking the event loop.
    """
    video_id_match = re.search(r"(?:v=|\/)([0-9A-Za-z_-]{11}).*", url)
    if not video_id_match:
        raise ValueError("Could not extract a valid YouTube Video ID from the provided URL.")

    video_id = video_id_match.group(1)

    try:
        # 1. Fetch transcript list
        transcripts = await asyncio.get_event_loop().run_in_executor(
            None,
            lambda: YouTubeTranscriptApi().list(video_id)
        )
        
        # 2. Try English first
        try:
            transcript_obj = transcripts.find_transcript(["en"])
        except Exception:
            # 3. Fallback to the first available transcript (e.g. Hindi, Spanish, or auto-generated)
            try:
                transcript_obj = next(iter(transcripts))
            except StopIteration:
                raise ValueError("No transcripts available for this video.")

        # 4. Fetch the transcript list content
        transcript = await asyncio.get_event_loop().run_in_executor(
            None,
            lambda: transcript_obj.fetch()
        )
    except Exception as e:
        if isinstance(e, ValueError):
            raise e
        raise ValueError(
            f"Could not retrieve transcript for video '{video_id}'. "
            "The video may not have closed captions or transcripts enabled."
        )

    extracted_chunks = []
    current_text_parts = []
    current_start_time = None
    current_length = 0

    for item in transcript:
        if current_start_time is None:
            current_start_time = item.start
        current_text_parts.append(item.text)
        current_length += len(item.text) + 1

        if current_length >= 800:
            combined_text = " ".join(current_text_parts)
            minutes = int(current_start_time // 60)
            seconds = int(current_start_time % 60)
            formatted_timestamp = f"{minutes}:{seconds:02d}"

            clean_text = combined_text.replace('\n', ' ').strip()
            if len(clean_text) > 10:
                extracted_chunks.append({
                    "content": clean_text,
                    "metadata": {
                        "source": url,
                        "timestamp": formatted_timestamp,
                        "type": "youtube"
                    }
                })

            current_text_parts = []
            current_start_time = None
            current_length = 0

    if current_text_parts:
        combined_text = " ".join(current_text_parts)
        minutes = int(current_start_time // 60)
        seconds = int(current_start_time % 60)
        formatted_timestamp = f"{minutes}:{seconds:02d}"

        clean_text = combined_text.replace('\n', ' ').strip()
        if len(clean_text) > 10:
            extracted_chunks.append({
                "content": clean_text,
                "metadata": {
                    "source": url,
                    "timestamp": formatted_timestamp,
                    "type": "youtube"
                }
            })

    return extracted_chunks


# ---------------------------------------------------------------------------
# Web Extraction  ||  Libraries: httpx + trafilatura + BeautifulSoup fallback
# ---------------------------------------------------------------------------
async def extract_text_from_web(url: str) -> List[Dict[str, Any]]:
    """
    Fetches a web page, extracts the main readable content, and prepares
    citation metadata for RAG ingestion.
    """
    parsed_url = urlparse(url)
    if parsed_url.scheme not in {"http", "https"} or not parsed_url.netloc:
        raise ValueError("A valid http or https URL must be provided for web extraction.")

    headers = {
        "User-Agent": (
            "Mozilla/5.0 (Windows NT 10.0; Win64; x64) "
            "AppleWebKit/537.36 (KHTML, like Gecko) "
            "Chrome/125.0 Safari/537.36"
        ),
        "Accept": "text/html,application/xhtml+xml,application/xml;q=0.9,*/*;q=0.8",
        "Accept-Language": "en-US,en;q=0.9",
    }

    try:
        async with httpx.AsyncClient(
            follow_redirects=True,
            timeout=httpx.Timeout(20.0, connect=10.0),
            headers=headers,
        ) as client:
            response = await client.get(url)
            response.raise_for_status()
    except httpx.HTTPStatusError as e:
        raise ValueError(
            f"Could not fetch web page '{url}'. Server returned status {e.response.status_code}."
        )
    except httpx.RequestError as e:
        raise ValueError(f"Could not fetch web page '{url}'. Error: {str(e)}")

    content_type = response.headers.get("content-type", "").lower()
    if "html" not in content_type and "xml" not in content_type:
        raise ValueError(f"The URL does not point to an HTML page. Content-Type: {content_type}")

    html = response.text
    title = _extract_html_title(html)
    full_markdown = trafilatura.extract(
        html,
        url=url,
        output_format="markdown",
        include_comments=False,
        include_tables=True,
        favor_precision=True,
    )

    if not full_markdown or len(full_markdown.strip()) < 100:
        full_markdown = _fallback_extract_text(html)

    if not full_markdown or len(full_markdown.strip()) < 50:
        raise ValueError("Could not extract enough readable text from the provided web page.")

    # Re-apply the chunking logic (~1000 chars) to match the expected return signature
    chunks = []
    current_chunk = ""
    metadata = {"source": url, "type": "web"}
    if title:
        metadata["title"] = title
    
    # Split by double newlines (standard paragraph breaks in Markdown)
    for paragraph in full_markdown.split("\n\n"):
        if not paragraph.strip():
            continue
            
        current_chunk += paragraph + "\n\n"
        
        if len(current_chunk) >= 1000:
            chunks.append({
                "content": current_chunk.strip(),
                "metadata": metadata
            })
            current_chunk = ""

    # Catch any remaining text at the end
    if current_chunk.strip():
        chunks.append({
            "content": current_chunk.strip(),
            "metadata": metadata
        })

    return chunks


def _extract_html_title(html: str) -> str | None:
    soup = BeautifulSoup(html, "lxml")
    if soup.title and soup.title.string:
        title = " ".join(soup.title.string.split())
        return title or None
    return None


def _fallback_extract_text(html: str) -> str:
    soup = BeautifulSoup(html, "lxml")
    for tag in soup(["script", "style", "noscript", "nav", "footer", "header", "form"]):
        tag.decompose()

    main = soup.find("main") or soup.find("article") or soup.body or soup
    paragraphs = [
        " ".join(element.get_text(" ", strip=True).split())
        for element in main.find_all(["h1", "h2", "h3", "p", "li", "blockquote"])
    ]
    return "\n\n".join(text for text in paragraphs if len(text) > 20)
